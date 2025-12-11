from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "University Capital / KEAN"
    subtitle.text = "Investment Case\nDecember 2025"
    
    # Slide 2: Executive Summary
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Executive Summary"
    tf = body_shape.text_frame
    
    tf.text = "Technology-Driven Platform with Strong IP"
    p = tf.add_paragraph()
    p.text = "• Current Valuation: $7.27M (weighted average)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 1.5M+ words across 99+ files"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 62,171+ technology mentions"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 2,062 development hours invested"
    p.level = 1
    
    # Slide 3: Investment Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Investment Highlights"
    
    # Add 3 content boxes
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(3)
    height = Inches(4)
    
    # Box 1: Tech Footprint
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "Tech Footprint"
    p.font.bold = True
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• 99+ code/docs\n• 1.5M+ words\n• 2,062 dev hours\n• 5x tech multiplier"
    
    # Box 2: Market & Valuation
    box = slide.shapes.add_textbox(left + width + 0.5, top, width, height)
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "Market & Valuation"
    p.font.bold = True
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• $7.27M valuation\n• $4.0M market approach\n• $18.0M income approach\n• 5.2x revenue multiplier"
    
    # Box 3: Growth Trajectory
    box = slide.shapes.add_textbox(left + (width + 0.5) * 2, top, width, height)
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = "Growth Trajectory"
    p.font.bold = True
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "• 1K → 5K → 25K users\n• $10K → $500K MRR\n• 15-300% value growth\n• 5-30 team size"
    
    # Slide 4: Market Positioning
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Market Positioning"
    
    # Add table
    rows = 4
    cols = 4
    left = Inches(1.0)
    top = Inches(1.5)
    width = Inches(8.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Table headers
    table.cell(0, 0).text = 'Metric'
    table.cell(0, 1).text = 'Competitor A'
    table.cell(0, 2).text = 'Competitor B'
    table.cell(0, 3).text = 'University Capital'
    
    # Table data
    table_data = [
        ['Valuation', '$5.0M', '$12.0M', '$7.27M'],
        ['Employees', '25', '50', '5 (FTE)'],
        ['Revenue', '$2.0M', '$5.0M', 'Projected $0.9M']
    ]
    
    for i, row in enumerate(table_data, 1):
        for j, value in enumerate(row):
            table.cell(i, j).text = value
    
    # Slide 5: Financial Projections
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Financial Projections"
    
    # Add table
    rows = 5
    cols = 4
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Table headers
    headers = ['Metric', 'Current', '1Y Projection', '3Y Projection']
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    # Table data
    table_data = [
        ['Users', '1,000', '5,000', '25,000'],
        ['MRR', '$10K', '$75K', '$500K'],
        ['Team Size', '5', '12', '30'],
        ['Valuation', '$7.27M', '$8.7-9.1M', '$10.9-12.7M']
    ]
    
    for i, row in enumerate(table_data, 1):
        for j, value in enumerate(row):
            table.cell(i, j).text = value
    
    # Slide 6: Investment Ask
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Investment Opportunity"
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    tf.text = "We are seeking [$X] in funding to:"
    p = tf.add_paragraph()
    p.text = "• Scale the development team"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Accelerate product development"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Expand market reach"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Enhance platform capabilities"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "\nProjected ROI: [X]% over [Y] years"
    
    # Save the presentation
    prs.save('UNIVERSITY_CAPITAL_CASE.pptx')
    print("Presentation generated successfully: UNIVERSITY_CAPITAL_CASE.pptx")

if __name__ == "__main__":
    create_presentation()
